"""
Build Canada Immigration & PR Navigator presentation — fixed layout v2.
Run: .venv/bin/python build_ppt.py
Output: presentation.pptx

Key fixes vs v1:
- Removed all line_spacing=Pt(X) calls (was setting FIXED pt spacing < font size → lines overlapping)
- Enlarged all text-box heights so content never clips outside bounds
- Reduced body font sizes to fit content comfortably
- Recalculated Y positions so nothing falls outside the 7.5" slide height
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
RED_ACCENT = RGBColor(0xC8, 0x10, 0x2E)
RED_DEEP   = RGBColor(0x8B, 0x1E, 0x2D)
TEXT_DARK  = RGBColor(0x1F, 0x1F, 0x1F)
BG_NEUTRAL = RGBColor(0xF6, 0xF3, 0xF1)
BORDER     = RGBColor(0xE5, 0xE0, 0xDC)
GRAY       = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0x99, 0x99, 0x99)
BG_DARK    = RGBColor(0x1F, 0x1F, 0x1F)

# ---------------------------------------------------------------------------
# Canvas: 13.33" × 7.5"
# ---------------------------------------------------------------------------
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

# ---------------------------------------------------------------------------
# Core helpers   (NO line_spacing override — default PowerPoint spacing)
# ---------------------------------------------------------------------------

def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s


def tb(slide, text, x, y, w, h,
       font="Calibri", size=Pt(16), bold=False, italic=False,
       color=TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True):
    """Single-run text box — no line_spacing override to avoid overlap."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font
    run.font.size   = size
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def blist(slide, items, x, y, w, h,
          font="Calibri", size=Pt(15), color=TEXT_DARK):
    """Bullet list with red em-dash, generous space_before."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(10)
        r1 = p.add_run(); r1.text = "— "
        r1.font.name = font; r1.font.size = size; r1.font.color.rgb = RED_ACCENT
        r2 = p.add_run(); r2.text = item
        r2.font.name = font; r2.font.size = size; r2.font.color.rgb = color
    return box


def title_bar(slide, title, tag, y=Inches(0.52)):
    tb(slide, title, Inches(0.8), y, Inches(9.5), Inches(0.6),
       font="Georgia", size=Pt(30), bold=True)
    if tag:
        tb(slide, tag.upper(), Inches(10.5), y + Inches(0.1), Inches(2.2), Inches(0.45),
           font="Calibri", size=Pt(11), bold=True, color=RED_ACCENT,
           align=PP_ALIGN.RIGHT, wrap=False)
    rect(slide, Inches(0.8), y + Inches(0.62), Inches(11.73), Pt(2), fill=TEXT_DARK)


# ===========================================================================
# SLIDE 1 — COVER
# ===========================================================================
def make_cover(prs):
    sl = prs.slides.add_slide(BLANK)
    rect(sl, Inches(0.8), Inches(0.75), Inches(1.1), Inches(0.08), fill=RED_ACCENT)

    # Title: 2 lines × Pt(50) ≈ ~1.6" min; give 2.0"
    tb(sl, "Canada Immigration\n& PR Navigator",
       Inches(0.8), Inches(1.0), Inches(10.5), Inches(2.0),
       font="Georgia", size=Pt(50), bold=True)

    # Subtitle: 2 lines × Pt(20) ≈ ~0.75"; give 1.05"
    tb(sl, "A RAG-powered policy assistant with citation-grounded,\nsafety-aware responses.",
       Inches(0.8), Inches(3.2), Inches(8.5), Inches(1.05),
       font="Calibri", size=Pt(20), color=GRAY)

    rect(sl, Inches(0.8), Inches(5.6), Inches(11.73), Pt(1), fill=BORDER)

    tb(sl, "TEAM 3",
       Inches(0.8), Inches(5.8), Inches(3.0), Inches(0.38),
       font="Calibri", size=Pt(13), bold=True)
    tb(sl, "Chao Tang  ·  Yuhan Ren  ·  Ehraaz Atif  ·  Ella Lu  ·  Keqing Wang",
       Inches(0.8), Inches(6.22), Inches(8.5), Inches(0.38),
       font="Calibri", size=Pt(14), color=GRAY)
    tb(sl, "April 2026",
       Inches(10.3), Inches(6.05), Inches(2.2), Inches(0.5),
       font="Georgia", size=Pt(18), italic=True, color=RED_DEEP,
       align=PP_ALIGN.RIGHT, wrap=False)


# ===========================================================================
# SLIDE 2 — EXECUTIVE SUMMARY
# ===========================================================================
def make_exec_summary(prs):
    sl = prs.slides.add_slide(BLANK)
    title_bar(sl, "Executive Summary", "Overview")

    # Red left-border accent bar
    rect(sl, Inches(0.8), Inches(1.62), Inches(0.055), Inches(0.92), fill=RED_ACCENT)

    # Impact statement: 2 lines × Pt(17) ≈ ~0.68"; give 1.0"
    tb(sl, "Delivering high-stakes immigration guidance with\nofficial-source traceability and robust risk mitigation.",
       Inches(1.0), Inches(1.58), Inches(6.0), Inches(1.0),
       font="Georgia", size=Pt(17), italic=True, color=RED_DEEP)

    # 3 bullets at Pt(16); each ~2 wrapped lines ≈ 0.55" + 10pt gap; 3 items ≈ 2.0"
    items = [
        "Developed an end-to-end MVP navigating complex IRCC and Ontario provincial policies.",
        "Implemented a strict citation schema ensuring responses are grounded in official documentation.",
        "Engineered a tiered safety-routing mechanism to manage adversarial or high-risk queries.",
    ]
    blist(sl, items, Inches(0.8), Inches(2.72), Inches(6.9), Inches(3.9), size=Pt(16))

    # KPI strip
    rx, ry, rw, rh = Inches(8.35), Inches(1.5), Inches(4.18), Inches(5.65)
    rect(sl, rx, ry, rw, rh, fill=BG_NEUTRAL)

    kpis = [
        ("72.73%", "EVAL PASS RATE (8 / 11)",     RED_ACCENT, Pt(44), "Georgia"),
        ("90.91%", "CITATION TITLE QUALITY",       TEXT_DARK,  Pt(44), "Georgia"),
        ("Hybrid", "RETRIEVAL: BM25 + VECTOR",     TEXT_DARK,  Pt(34), "Calibri"),
    ]
    ky = ry + Inches(0.22)
    for val, lbl, vc, vsz, vf in kpis:
        tb(sl, val, rx + Inches(0.25), ky, rw - Inches(0.5), Inches(0.82),
           font=vf, size=vsz, bold=True, color=vc)
        ky += Inches(0.8)
        tb(sl, lbl, rx + Inches(0.25), ky, rw - Inches(0.5), Inches(0.32),
           font="Calibri", size=Pt(11), bold=True, color=GRAY)
        ky += Inches(0.32)
        rect(sl, rx + Inches(0.25), ky, rw - Inches(0.5), Pt(1), fill=BORDER)
        ky += Inches(0.38)


# ===========================================================================
# SLIDE 3 — PROBLEM & DATASET
# ===========================================================================
def make_problem_dataset(prs):
    sl = prs.slides.add_slide(BLANK)
    title_bar(sl, "Problem & Dataset", "Foundation")

    tb(sl, "The Challenge",
       Inches(0.8), Inches(1.62), Inches(5.6), Inches(0.48),
       font="Georgia", size=Pt(21), bold=True)
    # body: Pt(15), 3 sentences ≈ 5–6 wrapped lines ≈ 1.65"; give 1.85"
    tb(sl, "Immigration guidance is exceptionally high-stakes, deeply policy-heavy, "
           "and subject to frequent legislative changes. Traditional LLM responses "
           "without strict grounding pose critical safety risks in this domain.",
       Inches(0.8), Inches(2.18), Inches(5.6), Inches(1.85),
       font="Calibri", size=Pt(15))

    rect(sl, Inches(6.82), Inches(1.62), Pt(1), Inches(2.1), fill=BORDER)

    tb(sl, "Source Authority",
       Inches(7.1), Inches(1.62), Inches(5.45), Inches(0.48),
       font="Georgia", size=Pt(21), bold=True)
    tb(sl, "Our dataset exclusively uses official IRCC and OINP webpages, processed into "
           "indexed chunks enriched with metadata: province, program, stream, source URL, "
           "section title, and accessed date — ensuring full traceability.",
       Inches(7.1), Inches(2.18), Inches(5.45), Inches(1.85),
       font="Calibri", size=Pt(15))

    # Flow diagram
    FY = Inches(4.25)
    FH = Inches(2.78)
    rect(sl, Inches(0.8), FY, Inches(11.73), FH, fill=BG_NEUTRAL)

    steps = [
        ("PHASE 01", "Raw Official\nPages",   "IRCC & OINP HTML"),
        ("PHASE 02", "Cleaned Text",          "Noise removal &\nnormalisation"),
        ("PHASE 03", "Semantic\nChunks",      "Context-aware\nsegmentation"),
        ("PHASE 04", "Indexed\nDatabase",     "Enriched with\nmetadata & vectors"),
    ]
    sw = Inches(2.32); arw = Inches(0.45); gap = Inches(0.32)
    sx0 = Inches(0.95); cy = FY + Inches(0.28); ch = Inches(2.24)

    for i, (ph, ttl, sub) in enumerate(steps):
        sx = sx0 + i * (sw + arw + gap)
        hi = (i == 3)
        rect(sl, sx, cy, sw, ch, fill=WHITE,
             line=RED_ACCENT if hi else BORDER, lw=Pt(1.5) if hi else Pt(0.75))
        if hi:
            rect(sl, sx, cy, sw, Pt(4), fill=RED_ACCENT)
        tb(sl, ph, sx + Inches(0.14), cy + Inches(0.18), sw - Inches(0.28), Inches(0.3),
           font="Calibri", size=Pt(10), bold=True, color=RED_ACCENT, wrap=False)
        tb(sl, ttl, sx + Inches(0.14), cy + Inches(0.52), sw - Inches(0.28), Inches(0.72),
           font="Calibri", size=Pt(15), bold=True)
        tb(sl, sub, sx + Inches(0.14), cy + Inches(1.3), sw - Inches(0.28), Inches(0.78),
           font="Calibri", size=Pt(13), color=GRAY)
        if i < 3:
            ax = sx + sw + gap * 0.28
            tb(sl, "→", ax, FY + Inches(1.42), arw + Inches(0.18), Inches(0.44),
               font="Calibri", size=Pt(22), bold=True, color=RED_ACCENT,
               align=PP_ALIGN.CENTER, wrap=False)


# ===========================================================================
# SLIDE 4 — ARCHITECTURE & DESIGN
# ===========================================================================
def make_architecture(prs):
    sl = prs.slides.add_slide(BLANK)
    title_bar(sl, "Architecture & Design", "System Pipeline")

    tb(sl, "Technical Decisions",
       Inches(0.8), Inches(1.65), Inches(6.0), Inches(0.52),
       font="Georgia", size=Pt(22), bold=True)

    blocks = [
        ("HYBRID RETRIEVAL STRATEGY",
         "BM25 (0.6) + Vector (0.4) weights to prioritize exact policy terms while maintaining semantic recall."),
        ("PRECISION FILTERING",
         "Initial top_k=20 reranked to final 5 with metadata filters: province, program, stream."),
        ("RISK ROUTING PROTOCOL",
         "Tiered L1/L2/L3 safety mechanism with one-retry no-evidence fallback to prevent hallucination."),
        ("POLICY TOOL INTEGRATION",
         "CRS calculator scoped to single-applicant Federal Express Entry for MVP delivery."),
    ]
    # 4 blocks × 1.08" = 4.32"; from y=2.28 reaches 6.60" ✓
    ty = Inches(2.28)
    for lbl, body in blocks:
        tb(sl, lbl, Inches(0.8), ty, Inches(5.9), Inches(0.3),
           font="Calibri", size=Pt(11), bold=True, color=RED_DEEP)
        # body: Pt(14), 1–2 wrapped lines ≈ 0.55"; give 0.68"
        tb(sl, body, Inches(0.8), ty + Inches(0.32), Inches(5.9), Inches(0.68),
           font="Calibri", size=Pt(14))
        ty += Inches(1.08)

    # Pipeline diagram
    PX = Inches(7.45); PW = Inches(5.1)
    rect(sl, PX, Inches(1.48), PW, Inches(5.72), fill=BG_NEUTRAL)
    tb(sl, "EXECUTION FLOW",
       PX + Inches(0.3), Inches(1.65), PW - Inches(0.6), Inches(0.3),
       font="Calibri", size=Pt(10), bold=True, color=GRAY)

    nodes = [
        ("1.  Intake Profile",         True,  False),
        ("2.  Hybrid Retrieval",        False, False),
        ("3.  Optional: Policy Tools",  False, False),
        ("4.  Tiered Risk Routing",     True,  False),
        ("5.  Generation + Citations",  False, True),
    ]
    # 5 nodes × 0.64" + 4 arrows × 0.20" = 3.20 + 0.80 = 4.00"; from 2.10 → 6.10" ✓
    ny = Inches(2.1); nh = Inches(0.64); aph = Inches(0.20)

    for text, hi, dark in nodes:
        bg = BG_DARK if dark else WHITE
        fg = WHITE  if dark else TEXT_DARK
        lc = RED_ACCENT if (hi or dark) else BORDER
        lw = Pt(2) if (hi or dark) else Pt(0.75)
        rect(sl, PX + Inches(0.3), ny, PW - Inches(0.6), nh, fill=bg, line=lc, lw=lw)
        if hi:
            rect(sl, PX + Inches(0.3), ny, Inches(0.06), nh, fill=RED_ACCENT)
        tb(sl, text, PX + Inches(0.48), ny + Inches(0.13),
           PW - Inches(0.85), Inches(0.42),
           font="Calibri", size=Pt(16), color=fg)
        ny += nh
        if text != "5.  Generation + Citations":
            tb(sl, "↓", PX + PW / 2 - Inches(0.18), ny,
               Inches(0.36), aph,
               font="Calibri", size=Pt(13), bold=True,
               color=RED_ACCENT, align=PP_ALIGN.CENTER, wrap=False)
            ny += aph


# ===========================================================================
# SLIDE 5 — EVALUATION RESULTS
# ===========================================================================
def make_eval_results(prs):
    sl = prs.slides.add_slide(BLANK)
    title_bar(sl, "Evaluation Results", "Metrics & Analysis")

    mx0 = Inches(0.8); mw = Inches(3.6); mg = Inches(0.165)

    metrics = [
        ("11 / 8",  "TOTAL TESTS / PASSED  (72.73%)", RED_ACCENT),
        ("90.91%",  "CITATION TITLE QUALITY",          TEXT_DARK),
        ("Valid",   "REFUSAL & SAFETY OUTCOMES",        TEXT_DARK),
    ]
    for i, (val, lbl, vc) in enumerate(metrics):
        bx = mx0 + i * (mw + mg)
        rect(sl, bx, Inches(1.62), mw, Inches(1.55), fill=WHITE, line=BORDER)
        tb(sl, val, bx + Inches(0.2), Inches(1.75), mw - Inches(0.4), Inches(0.85),
           font="Georgia", size=Pt(36), bold=True, color=vc)
        tb(sl, lbl, bx + Inches(0.2), Inches(2.58), mw - Inches(0.4), Inches(0.38),
           font="Calibri", size=Pt(11), bold=True, color=GRAY)

    # Failure modes
    tb(sl, "Observed Failure Modes",
       Inches(0.8), Inches(3.35), Inches(6.5), Inches(0.44),
       font="Calibri", size=Pt(19), bold=True)
    fails = [
        "L3 Compliance Inconsistency: Variable refusal behavior across diverse adversarial safety prompts.",
        "Factual Keyword Misses: Retrieval gaps for highly specific, long-tail policy detail queries.",
        "Sample Size Constraints: Small evaluation set limits confidence around broad generalization.",
    ]
    # 3 bullets × ~0.85" each ≈ 2.55"; give 3.1"
    blist(sl, fails, Inches(0.8), Inches(3.88), Inches(6.6), Inches(3.1), size=Pt(15))

    # Dark interpretation block
    bx, by_, bw, bh = Inches(7.7), Inches(3.28), Inches(4.85), Inches(3.85)
    rect(sl, bx, by_, bw, bh, fill=BG_DARK)
    tb(sl, "WHAT THIS MEANS",
       bx + Inches(0.3), by_ + Inches(0.25), bw - Inches(0.6), Inches(0.32),
       font="Calibri", size=Pt(11), bold=True, color=GRAY_LIGHT)
    # body: Pt(15), ~5 wrapped lines ≈ 1.65"; give 2.95"
    tb(sl, "The system demonstrates strong capability in tracing and citing official "
           "sources accurately, meeting our primary traceability mandate. Edge-case "
           "safety routing and granular retrieval tuning still require further dataset "
           "expansion before production deployment.",
       bx + Inches(0.3), by_ + Inches(0.68), bw - Inches(0.6), Inches(2.95),
       font="Calibri", size=Pt(15), color=WHITE)


# ===========================================================================
# SLIDE 6 — LEARNINGS & NEXT STEPS
# ===========================================================================
def make_learnings(prs):
    sl = prs.slides.add_slide(BLANK)
    title_bar(sl, "Learnings & Next Steps", "Project Retrospective")

    tiles = [
        {
            "title": "What Worked Well",
            "top": RED_ACCENT, "bg": BG_NEUTRAL,
            "items": [
                "Hybrid retrieval significantly improved policy grounding over pure semantic search.",
                "Strict citation schemas enforced end-to-end response traceability.",
                "Modular orchestrator enabled efficient parallel team development.",
            ],
        },
        {
            "title": "Key Challenges",
            "top": TEXT_DARK, "bg": BG_NEUTRAL,
            "items": [
                "Consistent refusal behavior under complex adversarial prompt injections.",
                "Balancing conversational personalization against evidence sufficiency requirements.",
                "Data freshness and robust ingestion under tight MVP time constraints.",
            ],
        },
        {
            "title": "Next Steps",
            "top": BORDER, "bg": WHITE,
            "items": [
                "Expand evaluation set significantly beyond the initial 11 seeds.",
                "Strengthen L3 refusal reliability and automated citation quality checks.",
                "Extend policy tool coverage beyond single-applicant Federal Express Entry.",
            ],
        },
    ]

    tw = Inches(3.72); tg = Inches(0.175)
    tx = Inches(0.8);  ty = Inches(1.62); th = Inches(5.52)

    for i, tile in enumerate(tiles):
        bx = tx + i * (tw + tg)
        rect(sl, bx, ty, tw, th,
             fill=tile["bg"],
             line=BORDER if tile["bg"] == WHITE else None)
        rect(sl, bx, ty, tw, Inches(0.07), fill=tile["top"])
        tb(sl, tile["title"],
           bx + Inches(0.25), ty + Inches(0.18), tw - Inches(0.5), Inches(0.52),
           font="Georgia", size=Pt(21), bold=True)
        # 3 bullets × ~0.8" each ≈ 2.4"; give 4.4"
        blist(sl, tile["items"],
              bx + Inches(0.15), ty + Inches(0.82), tw - Inches(0.3), Inches(4.4),
              size=Pt(14))


# ===========================================================================
# SLIDE 7 — THANK YOU
# ===========================================================================
def make_thankyou(prs):
    sl = prs.slides.add_slide(BLANK)
    rect(sl, Inches(5.77), Inches(1.95), Inches(1.8), Inches(0.07), fill=RED_ACCENT)
    tb(sl, "Thank You.",
       Inches(0.8), Inches(2.2), Inches(11.73), Inches(1.65),
       font="Georgia", size=Pt(68), bold=True,
       align=PP_ALIGN.CENTER)
    tb(sl, "Questions welcome.",
       Inches(0.8), Inches(4.1), Inches(11.73), Inches(0.65),
       font="Georgia", size=Pt(26), italic=True,
       color=RED_DEEP, align=PP_ALIGN.CENTER)
    tb(sl, "CANADA IMMIGRATION & PR NAVIGATOR  ·  RAG MVP  ·  TEAM 3",
       Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.42),
       font="Calibri", size=Pt(12),
       color=GRAY_LIGHT, align=PP_ALIGN.CENTER)


# ===========================================================================
# BUILD
# ===========================================================================
make_cover(prs)
make_exec_summary(prs)
make_problem_dataset(prs)
make_architecture(prs)
make_eval_results(prs)
make_learnings(prs)
make_thankyou(prs)

OUTPUT = "presentation.pptx"
prs.save(OUTPUT)
print(f"Saved → {OUTPUT}  ({round(__import__('os').path.getsize(OUTPUT)/1024)} KB)")
